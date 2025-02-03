import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image


def pdf_to_ppt(pdf_path, ppt_path):
    # Open the PDF file
    doc = fitz.open(pdf_path)
    prs = Presentation()

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)

        # Extract text
        text = page.get_text("text")

        # Create a new slide
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        text_frame.text = text

        # Extract images
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            # Convert image to Pillow format
            image = Image.open(io.BytesIO(image_bytes))

            # Save image to a temporary file
            img_path = f"temp_image_{page_num + 1}_{img_index + 1}.{image_ext}"
            image.save(img_path)

            # Add image to slide
            slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(6))

    # Save the PowerPoint file
    prs.save(ppt_path)
    print(f"PowerPoint file saved: {ppt_path}")


# Example Usage
pdf_to_ppt("sample.pdf", "output.pptx")
