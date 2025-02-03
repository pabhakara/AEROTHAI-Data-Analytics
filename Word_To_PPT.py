from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2 import service_account
from pptx import Presentation
import io

# Google Drive API Setup
SCOPES = ['https://www.googleapis.com/auth/drive']
SERVICE_ACCOUNT_FILE = 'credentials.json'  # Update with your file path


def authenticate_google_drive():
    """Authenticate using Google Service Account"""
    creds = service_account.Credentials.from_service_account_file(
        SERVICE_ACCOUNT_FILE, scopes=SCOPES
    )
    return build('drive', 'v3', credentials=creds)


def download_google_doc_as_text(drive_service, file_id):
    """Download a Google Doc as plain text"""
    request = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
    file_data = io.BytesIO()
    downloader = MediaIoBaseDownload(file_data, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()

    return file_data.getvalue().decode('utf-8')


def create_ppt_from_text(text, output_pptx):
    """Convert extracted text into PowerPoint slides"""
    prs = Presentation()

    for para in text.split("\n\n"):  # Split by double newlines (paragraphs)
        if para.strip():  # Avoid empty slides
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
            title, content = para.split("\n", 1) if "\n" in para else (para, "")
            slide.shapes.title.text = title
            slide.shapes.placeholders[1].text = content  # Content box

    prs.save(output_pptx)
    print(f"PowerPoint saved as: {output_pptx}")


def main():
    drive_service = authenticate_google_drive()

    GOOGLE_DOC_ID = 'your_google_doc_id_here'  # Replace with your Google Doc ID
    OUTPUT_PPTX = 'output_presentation.pptx'

    # Download the Google Doc as text
    text_content = download_google_doc_as_text(drive_service, GOOGLE_DOC_ID)

    # Convert text to PowerPoint
    create_ppt_from_text(text_content, OUTPUT_PPTX)


if __name__ == "__main__":
    main()
